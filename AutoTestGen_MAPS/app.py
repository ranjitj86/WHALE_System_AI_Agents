from flask import Flask, render_template, request, jsonify, send_file, make_response, session, Response
from dotenv import load_dotenv
import os
import webbrowser
from agents.intake_agent import ElicitationAgent
from agents.sys2_agent import Sys2Agent
from agents.review_agent import ReviewAgent
from agents.testgen_agent import TestGenAgent
import pptx
import email
import extract_msg
import pandas as pd
from typing import List, Dict, Any
import io
import csv
import docx as docx_lib
from fpdf import FPDF
import time
import threading
from datetime import datetime

# Load environment variables
load_dotenv()

app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'your-secret-key-here')
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Ensure upload directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Initialize agents
elicitation_agent = ElicitationAgent()
sys2_agent = Sys2Agent()
review_agent = ReviewAgent()
testgen_agent = TestGenAgent() # Initialize Agent 4

# In-memory cache for last extracted requirements (REMOVING - using session instead)
# last_elicitation_requirements = []

@app.route('/')
def index():
    """Landing page with agent selector interface"""
    # Ensure session is initialized if not already
    if 'activity_messages' not in session:
        session['activity_messages'] = []

    # Get relevant status data from session for display on the landing page cards
    sys1_requirements = session.get('sys1_elicitation_requirements', [])

    # Calculate status for Agent 1 (Elicitation)
    total_sys1_reqs = len(sys1_requirements)
    approved_sys1_count = sum(1 for req in sys1_requirements if req.get('req_status') == 'Approved')
    rejected_sys1_count = sum(1 for req in sys1_requirements if req.get('req_status') == 'Rejected')
    draft_sys1_count = total_sys1_reqs - approved_sys1_count - rejected_sys1_count

    agent1_status = f"{total_sys1_reqs} requirements extracted. ({approved_sys1_count} Approved, {rejected_sys1_count} Rejected, {draft_sys1_count} Draft)" if total_sys1_reqs > 0 else "Ready for inputs."

    # Placeholder status for other agents (can be made dynamic later)
    agent2_status = "Under Development" # Check for SYS.2 draft data in session if available
    agent3_status = f"{sum(1 for req in sys1_requirements if req.get('req_status') == 'Draft')} items need refinement" if draft_sys1_count > 0 else "Ready for review"
    agent4_status = "Under Development" # Check for generated test cases in session if available

    return render_template('index.html', 
                           agent1_status=agent1_status,
                           agent2_status=agent2_status,
                           agent3_status=agent3_status,
                           agent4_status=agent4_status,
                           )

@app.route('/agent/<agent_id>')
def agent_dashboard(agent_id):
    """Dynamic agent dashboard routing"""
    if agent_id == '1':
        # Load SYS.1 requirements from session for Agent 1 dashboard display
        sys1_requirements = session.get('sys1_elicitation_requirements', [])
        print(f"[DEBUG] SYS.1 Requirements loaded from session for Agent 1: {len(sys1_requirements)} requirements") # Debug print
        # Pass only SYS.1 requirements to the template
        return render_template('agent1/dashboard.html', initial_requirements=sys1_requirements)
    elif agent_id == '4':
        # Pass necessary data for Agent 4 dashboard
        # We might load requirements automatically here later, but for now, just render the template
        # Pass the default input path to the template
        return render_template('agent4/dashboard.html', 
                               agent_name=testgen_agent.agent_name,
                               agent_description=testgen_agent.agent_description,
                               default_input_path=testgen_agent.default_input_path,
                               # Pass any existing loaded requirements or generated test cases from session
                               loaded_requirements=session.get('agent4_requirements', []),
                               generated_test_cases=session.get('agent4_test_cases', [])
                              )
    else:
        # Other agents can be handled here if they need initial data loaded
        return render_template(f'agent{agent_id}/dashboard.html')

@app.route('/traceability')
def traceability_dashboard():
    return render_template('traceability_dashboard/dashboard.html')

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Handle file uploads for all agents and extract SYS.1 requirements for Agent 1"""
    # global last_elicitation_requirements # REMOVING
    files = request.files.getlist('file')
    raw_content = request.form.get('raw_content', '')
    all_text = ''
    for file in files:
        if file.filename == '':
            continue
        filename = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filename)
        ext = file.filename.split('.')[-1].lower()
        try:
            if ext == 'txt':
                with open(filename, 'r', encoding='utf-8', errors='ignore') as f:
                    all_text += f.read() + '\n'
            elif ext == 'pdf':
                import PyPDF2
                with open(filename, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        all_text += page.extract_text() + '\n'
            elif ext == 'docx':
                import docx
                doc = docx.Document(filename)
                for para in doc.paragraphs:
                    all_text += para.text + '\n'
            elif ext in ['xlsx', 'csv']:
                df = pd.read_excel(filename) if ext == 'xlsx' else pd.read_csv(filename)
                all_text += df.to_string(index=False) + '\n'
            elif ext in ['ppt', 'pptx']:
                prs = pptx.Presentation(filename)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            all_text += shape.text + '\n'
            elif ext == 'eml':
                msg = email.message_from_file(open(filename, 'r', encoding='utf-8', errors='ignore'))
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        all_text += part.get_payload(decode=True).decode(errors='ignore') + '\n'
            elif ext == 'msg':
                msg = extract_msg.Message(filename)
                all_text += msg.body + '\n'
            # Add more handlers for Jira/Confluence as needed
            else:
                with open(filename, 'r', encoding='utf-8', errors='ignore') as f:
                    all_text += f.read() + '\n'
        except Exception as e:
            all_text += f'\n[Error reading {file.filename}: {e}]\n'
    if raw_content:
        all_text += raw_content + '\n'
    try:
        result = elicitation_agent.process({'content': all_text, 'format': 'multi'})
        if result.get('status') == 'success':
            customer_reqs = result.get('customer_requirements', [])
            sys1_reqs = result.get('sys1_requirements', [])

            # Store both lists in the session
            session['customer_elicitation_requirements'] = customer_reqs
            session['sys1_elicitation_requirements'] = sys1_reqs

            print(f"[DEBUG] Customer Requirements saved to session: {len(customer_reqs)} requirements") # Debug print
            print(f"[DEBUG] SYS.1 Requirements saved to session: {len(sys1_reqs)} requirements") # Debug print

            # --- Automatic Export of SYS.1 Requirements Only to XLSX ----
            try:
                # Use the sys1_reqs directly, no need to get from session again immediately
                sys1_only_data = [{
                    'SYS.1 Req. ID': req.get('sys1_id', ''),
                    'SYS.1 System Requirement': req.get('sys1_requirement', '')
                } for req in sys1_reqs]

                if sys1_only_data:
                    df = pd.DataFrame(sys1_only_data)

                    # Define the path to save the file in the root directory
                    # Get the directory where app.py is located (assuming it's the root for now)
                    root_path = os.path.dirname(os.path.abspath(__file__))
                    # Define the path to save the file in the Inputs directory
                    inputs_dir = os.path.join(root_path, 'Inputs')
                    os.makedirs(inputs_dir, exist_ok=True) # Create Inputs directory if it doesn't exist
                    excel_file_path = os.path.join(inputs_dir, 'sys1_requirements.xlsx')

                    df.to_excel(excel_file_path, index=False)
                    print(f"[INFO] Successfully exported SYS.1 requirements to {excel_file_path}")
                    file_export_message = f'SYS.1 requirements exported to {excel_file_path}'

            except Exception as e:
                print(f"[ERROR] Failed to automatically export SYS.1 requirements: {e}")
                # Add an error message to activity feed (optional)
                # if 'activity_messages' in session:
                #     session['activity_messages'].append(f'Failed to export SYS.1 requirements: {e}')
                #     session.modified = True
            # --- End Automatic Export ---

            # Return SYS.1 requirements for immediate display on Agent 1 page
            response_data = {
                'status': 'success',
                'requirements': sys1_reqs,
                'metadata': result.get('metadata', {})
            }
            # Add the file export message if it was set
            if 'file_export_message' in locals():
                 response_data['file_export_message'] = file_export_message
            elif 'excel_file_path' in locals():
                 # If export path was determined but failed, maybe provide a generic message
                 response_data['file_export_message'] = f'Attempted to export SYS.1 requirements to {excel_file_path}, check terminal for status.'

            return jsonify(response_data)
        else:
            return jsonify({'status': 'error', 'message': result.get('message', 'Unknown error')}), 500
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

# Agent 3 (Review) endpoints
@app.route('/api/agent3/review', methods=['POST'])
def review_requirements():
    try:
        data = request.get_json()
        if not data:
            return jsonify({'status': 'error', 'message': 'No data provided'}), 400

        # Process the requirements for review
        result = review_agent.process(data)
        
        return jsonify({
            'status': 'success',
            'requirements': result.get('requirements', []),
            'compliance_results': result.get('compliance_results', []),
            'suggestions': result.get('suggestions', []),
            'test_proposals': result.get('test_proposals', [])
        })
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/agent3/apply-suggestion', methods=['POST'])
def apply_suggestion():
    try:
        data = request.get_json()
        if not data or 'requirement_id' not in data or 'suggestion' not in data:
            return jsonify({'status': 'error', 'message': 'Missing required fields'}), 400

        result = review_agent.apply_suggestion(data['requirement_id'], data['suggestion'])
        
        return jsonify({
            'status': 'success',
            'updated_requirement': result
        })
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/agent3/update-compliance', methods=['POST'])
def update_compliance():
    try:
        data = request.get_json()
        if not data or 'requirement_id' not in data or 'status' not in data:
            return jsonify({'status': 'error', 'message': 'Missing required fields'}), 400

        result = review_agent.update_compliance_status(data['requirement_id'], data['status'])
        
        return jsonify({
            'status': 'success',
            'updated_status': result
        })
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/all_agents')
def all_agents_page():
    return render_template('all_agents.html')

@app.route('/api/all_agents/process', methods=['POST'])
def all_agents_process():
    try:
        # 1. Save uploaded file
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'}), 400
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No selected file'}), 400
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filepath)

        results = {}
        # 2. Agent 1: Elicitation (SYS.1)
        try:
            sys1_result = elicitation_agent.process({'file_path': filepath})
            results['sys1'] = {'result': str(sys1_result)}
        except Exception as e:
            results['sys1'] = {'error': str(e)}

        # 3. Agent 2: SYS.2 Drafting
        try:
            sys2_result = sys2_agent.process({'requirements': sys1_result})
            results['sys2'] = {'result': str(sys2_result)}
        except Exception as e:
            results['sys2'] = {'error': str(e)}

        # 4. Agent 3: Review
        try:
            review_result = review_agent.process({'requirements': sys2_result})
            results['review'] = {'result': str(review_result)}
        except Exception as e:
            results['review'] = {'error': str(e)}

        # 5. Agent 2: SYS.2 Final Drafting (if needed, can be same as previous)
        try:
            sys2_final_result = sys2_agent.process({'requirements': review_result})
            results['sys2_final'] = {'result': str(sys2_final_result)}
        except Exception as e:
            results['sys2_final'] = {'error': str(e)}

        # 6. Agent 4: SYS.5 Test Case Generation
        try:
            sys5_result = testgen_agent.process({'requirements': sys2_final_result})
            results['sys5'] = {'result': str(sys5_result)}
        except Exception as e:
            results['sys5'] = {'error': str(e)}

        return jsonify(results)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/agent1/export/<format>', methods=['GET'])
def export_requirements(format):
    # Load SYS.1 requirements from session for export
    sys1_requirements_to_export = session.get('sys1_elicitation_requirements', [])
    customer_requirements_for_trace = session.get('customer_elicitation_requirements', [])

    # Create a map of customer requirements for easy lookup
    customer_map = {req.get('customer_id'): req for req in customer_requirements_for_trace if req.get('customer_id')}

    # Calculate Summary Statistics
    total_customer_reqs = len(customer_requirements_for_trace)
    traced_customer_ids = set()
    for sys1_req in sys1_requirements_to_export:
        if sys1_req.get('customer_trace_ids'):
            for cust_id in sys1_req['customer_trace_ids']:
                if cust_id in customer_map: # Ensure the traced ID is valid
                    traced_customer_ids.add(cust_id)

    traced_customer_count = len(traced_customer_ids)
    untraced_customer_count = total_customer_reqs - traced_customer_count

    total_sys1_reqs = len(sys1_requirements_to_export)
    approved_sys1_count = sum(1 for req in sys1_requirements_to_export if req.get('req_status') == 'Approved')
    rejected_sys1_count = sum(1 for req in sys1_requirements_to_export if req.get('req_status') == 'Rejected')
    draft_sys1_count = total_sys1_reqs - approved_sys1_count - rejected_sys1_count

    summary_data = {
        'customer_traceability': {
            'total': total_customer_reqs,
            'traced': traced_customer_count,
            'untraced': untraced_customer_count
        },
        'sys1_status': {
            'total': total_sys1_reqs,
            'approved': approved_sys1_count,
            'rejected': rejected_sys1_count,
            'draft': draft_sys1_count
        }
    }

    if not sys1_requirements_to_export:
        return jsonify({'status': 'error', 'message': 'No requirements to export'}), 400

    if format == 'csv':
        output = io.StringIO()
        writer = csv.writer(output)
        # Add Summary Section to CSV
        writer.writerow(['Summary:'])
        writer.writerow(['Customer Traceability', '', 'SYS.1 Status'])
        writer.writerow(['Total Customer Req.', summary_data['customer_traceability']['total'], 'Total SYS.1 Req.', summary_data['sys1_status']['total'] ])
        writer.writerow(['Traced to SYS.1', summary_data['customer_traceability']['traced'], 'SYS.1 Approved', summary_data['sys1_status']['approved'] ])
        writer.writerow(['Not Traced to SYS.1', summary_data['customer_traceability']['untraced'], 'SYS.1 Rejected', summary_data['sys1_status']['rejected'] ])
        writer.writerow(['', '', 'SYS.1 Draft', summary_data['sys1_status']['draft'] ])
        writer.writerow([]) # Add an empty row for separation

        # Add Table Headers to CSV
        writer.writerow(['Customer Req. ID(s)', 'Customer Requirement', 'SYS.1 Req. ID', 'SYS.1 System Requirement', 'Domain', 'Priority', 'Rationale', 'Requirement Status'])
        for req in sys1_requirements_to_export:
            # Get linked customer requirement text (joining if multiple traces exist)
            customer_req_texts = [
                customer_map.get(cid, {}).get('customer_requirement', '')
                for cid in req.get('customer_trace_ids', [])
            ]
            writer.writerow([
                ', '.join(req.get('customer_trace_ids', [])),
                '; '.join(customer_req_texts), # Use semicolon to separate texts if multiple
                req.get('sys1_id', ''),
                req.get('sys1_requirement', ''),
                req.get('domain', ''),
                req.get('priority', ''),
                req.get('rationale', ''),
                req.get('req_status', 'Draft')
            ])
        output.seek(0)
        return send_file(io.BytesIO(output.getvalue().encode('utf-8')), mimetype='text/csv', as_attachment=True, download_name='elicitation_requirements.csv')
    elif format == 'xlsx':
        import pandas as pd
        # Prepare data for DataFrame, converting list of IDs to string
        # Add Customer Requirement columns to DataFrame data
        df_data = []
        for req in sys1_requirements_to_export:
            req_copy = req.copy()
            req_copy['Customer Req. ID(s)'] = ', '.join(req_copy.get('customer_trace_ids', []))
            # Get linked customer requirement text
            customer_req_texts = [
                customer_map.get(cid, {}).get('customer_requirement', '') for cid in req.get('customer_trace_ids', [])
            ]
            req_copy['Customer Requirement'] = '; '.join(customer_req_texts)
            # Rename existing keys to match desired output headers if necessary, or add new ones
            req_copy['SYS.1 Req. ID'] = req_copy.pop('sys1_id', '')
            req_copy['SYS.1 System Requirement'] = req_copy.pop('sys1_requirement', '')
            req_copy['Domain'] = req_copy.pop('domain', '')
            req_copy['Priority'] = req_copy.pop('priority', '')
            req_copy['Rationale'] = req_copy.pop('rationale', '')
            req_copy['Requirement Status'] = req_copy.pop('req_status', '')
            # Remove internal keys not needed in export
            req_copy.pop('customer_trace_ids', None)
            df_data.append(req_copy)
        # Debugging: Print df_data before creating DataFrame
        print("[DEBUG] Data for XLSX DataFrame:", df_data)
        df = pd.DataFrame(df_data)
        # Reorder columns to match desired export format (Customer Req. first)
        # Temporarily simplify columns for debugging XLSX export issue
        cols = ['SYS.1 Req. ID', 'SYS.1 System Requirement'] # Simplified columns for testing
        df = df.reindex(columns=cols)
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            # Add Summary Sheet to XLSX
            summary_df = pd.DataFrame({
                'Summary Type': ['Customer Traceability', '', 'SYS.1 Status', '', ''],
                'Metric': ['Total Customer Req.', 'Traced to SYS.1', 'Not Traced to SYS.1', 'Total SYS.1 Req.', 'SYS.1 Approved', 'SYS.1 Rejected', 'SYS.1 Draft'],
                'Count': [
                    summary_data['customer_traceability']['total'],
                    summary_data['customer_traceability']['traced'],
                    summary_data['customer_traceability']['untraced'],
                    summary_data['sys1_status']['total'],
                    summary_data['sys1_status']['approved'],
                    summary_data['sys1_status']['rejected'],
                    summary_data['sys1_status']['draft']
                ]
            })
            summary_df.to_excel(writer, sheet_name='Summary', index=False)

            df.to_excel(writer, index=False, sheet_name='SYS1 Requirements')
        output.seek(0)
        # Ensure the BytesIO object is ready to be read from the beginning
        output.seek(0) # Add this line to reset the stream position to the beginning
        return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', as_attachment=True, download_name='elicitation_requirements.xlsx')
    elif format == 'docx':
        doc = docx_lib.Document()
        # Add Summary Section to DOCX
        doc.add_heading('Summary', level=1)
        doc.add_paragraph(f"Total Customer Requirements: {summary_data['customer_traceability']['total']}")
        doc.add_paragraph(f"Customer Requirements Traced to SYS.1: {summary_data['customer_traceability']['traced']}")
        doc.add_paragraph(f"Customer Requirements Not Traced to SYS.1: {summary_data['customer_traceability']['untraced']}")
        doc.add_paragraph(f"Total SYS.1 Requirements: {summary_data['sys1_status']['total']}")
        doc.add_paragraph(f"SYS.1 Requirements Approved: {summary_data['sys1_status']['approved']}")
        doc.add_paragraph(f"SYS.1 Requirements Rejected: {summary_data['sys1_status']['rejected']}")
        doc.add_paragraph(f"SYS.1 Requirements Draft: {summary_data['sys1_status']['draft']}")

        doc.add_heading('SYS.1 Requirements', level=1)
        # Add Table Headers to DOCX
        table = doc.add_table(rows=1, cols=8) # 8 columns now
        table.style = 'Table Grid' # Apply a grid style for better visibility
        hdr_cells = table.rows[0].cells
        headers = ['Customer Req. ID(s)', 'Customer Requirement', 'SYS.1 Req. ID', 'SYS.1 System Requirement', 'Domain', 'Priority', 'Rationale', 'Requirement Status']
        for i, header in enumerate(headers):
            hdr_cells[i].text = header

        for req in sys1_requirements_to_export:
            row_cells = table.add_row().cells
            row_cells[0].text = ', '.join(req.get('customer_trace_ids', []))
            row_cells[1].text = '; '.join(customer_req_texts)
            row_cells[2].text = req.get('sys1_id', '')
            row_cells[3].text = req.get('sys1_requirement', '')
            row_cells[4].text = req.get('domain', '')
            row_cells[5].text = req.get('priority', '')
            row_cells[6].text = req.get('rationale', '')
            row_cells[7].text = req.get('req_status', 'Draft')

        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document', as_attachment=True, download_name='elicitation_requirements.docx')
    elif format == 'pdf':
        class PDF(FPDF):
            def header(self):
                self.set_font('Arial', 'B', 12)
                self.cell(0, 10, 'Requirement Elicitation Report', 0, 1, 'C')
                self.ln(5)

            def footer(self):
                self.set_y(-15)
                self.set_font('Arial', 'I', 8)
                self.cell(0, 10, f'Page {self.page_no()}/{{nb}}', 0, 0, 'C')

        pdf = PDF()
        pdf.alias_nb_pages()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Add heading for Extracted Requirements
        pdf.set_font('Arial', 'B', 12)
        pdf.cell(0, 10, 'Extracted Requirements', 0, 1, 'L')
        pdf.ln(2)

        # Table Section
        headers = ['Customer Req. ID(s)', 'Customer Req.', 'SYS.1 Req. ID', 'SYS.1 System Req.', 'Domain', 'Priority', 'Rationale', 'Status']
        col_widths = [30, 50, 20, 50, 20, 20, 40, 20]

        pdf.set_font('Arial', 'B', 10)
        for i, header in enumerate(headers):
            pdf.cell(col_widths[i], 10, header, 1, 0, 'C')
        pdf.ln()

        pdf.set_font('Arial', '', 8)
        for req in sys1_requirements_to_export:
            customer_req_texts = [
                customer_map.get(cid, {}).get('customer_requirement', '') for cid in req.get('customer_trace_ids', [])
            ]
            cell_data = [
                ', '.join(req.get('customer_trace_ids', [])),
                '; '.join(customer_req_texts),
                req.get('sys1_id', ''),
                req.get('sys1_requirement', ''),
                req.get('domain', ''),
                req.get('priority', ''),
                req.get('rationale', ''),
                req.get('req_status', 'Draft')
            ]
            for i, data in enumerate(cell_data):
                pdf.cell(col_widths[i], 10, str(data), 1, 0, 'L')
            pdf.ln()

        # Add some space before the dashboard summary
        pdf.ln(8)
        pdf.set_font('Arial', 'B', 12)
        pdf.cell(0, 10, 'Traceability and Status Dashboard', 0, 1, 'L')
        pdf.ln(2)

        # Print summary in table format
        pdf.set_font('Arial', 'B', 10)
        pdf.cell(60, 8, 'Metric', 1, 0, 'C')
        pdf.cell(40, 8, 'Value', 1, 1, 'C')
        pdf.set_font('Arial', '', 10)
        summary_rows = [
            ('Total Customer Requirements', summary_data['customer_traceability']['total']),
            ('Customer Requirements Traced to SYS.1', summary_data['customer_traceability']['traced']),
            ('Customer Requirements Not Traced to SYS.1', summary_data['customer_traceability']['untraced']),
            ('Total SYS.1 Requirements', summary_data['sys1_status']['total']),
            ('SYS.1 Requirements Approved', summary_data['sys1_status']['approved']),
            ('SYS.1 Requirements Rejected', summary_data['sys1_status']['rejected']),
            ('SYS.1 Requirements Draft', summary_data['sys1_status']['draft'])
        ]
        for metric, value in summary_rows:
            pdf.cell(60, 8, metric, 1, 0, 'L')
            pdf.cell(40, 8, str(value), 1, 1, 'C')

        output = io.BytesIO()
        # Use dest='B' to get the raw bytes output
        pdf_bytes = pdf.output(dest='B')
        output.write(pdf_bytes)
        output.seek(0)
        return send_file(output, mimetype='application/pdf', as_attachment=True, download_name='elicitation_requirements.pdf')
    else:
        return jsonify({'status': 'error', 'message': 'Unsupported export format'}), 400

@app.route('/api/agent1/export/sys1_only_xlsx', methods=['GET'])
def export_sys1_only_xlsx():
    requirements_data = session.get('sys1_elicitation_requirements', [])

    if not requirements_data:
        return jsonify({'status': 'error', 'message': 'No requirements to export.'}), 400

    # Extract only SYS.1 System Requirement and SYS.1 Req. ID (for context/identifier)
    sys1_only_data = [{
        'SYS.1 Req. ID': req.get('sys1_id', ''),
        'SYS.1 System Requirement': req.get('sys1_requirement', '')
    } for req in requirements_data]

    df = pd.DataFrame(sys1_only_data)

    # Define the path to save the file in the root directory
    # The workspace root is available from the user_info. Assuming relative path from app.py location
    # This might need adjustment based on actual app.py location relative to project root.
    # Let's assume app.py is in the root or we can derive the root path.
    # Using a hardcoded path derived from user_info for now. C:\Users\Ranjit Jagtap\Desktop\AutoTestGen_Project\AutoTestGen_MAPS
    # A more robust approach might involve a configuration or calculating the path.
    root_path = os.path.dirname(os.path.abspath(__file__))
    excel_file_path = os.path.join(root_path, 'sys1_requirements_only.xlsx')

    try:
        df.to_excel(excel_file_path, index=False)
        return jsonify({'status': 'success', 'message': f'SYS.1 requirements exported to {excel_file_path}'})
    except Exception as e:
        print(f"Error exporting SYS.1 requirements: {e}")
        return jsonify({'status': 'error', 'message': f'Failed to export SYS.1 requirements: {str(e)}'}), 500

@app.route('/api/agent1/update_status', methods=['POST'])
def update_requirement_status():
    # We now update SYS.1 requirements based on sys1_id
    data = request.json
    sys1_id = data.get('id') # Use sys1_id now
    new_req_status = data.get('new_req_status')

    if not sys1_id or not new_req_status:
        return jsonify({'status': 'error', 'message': 'Missing id or new_req_status'}), 400

    # Load SYS.1 requirements from session
    sys1_requirements = session.get('sys1_elicitation_requirements', [])
    for req in sys1_requirements:
        if req.get('sys1_id') == sys1_id:
            req['req_status'] = new_req_status
            session['sys1_elicitation_requirements'] = sys1_requirements # Save updated list back to session
            return jsonify({'status': 'success', 'message': f'Status updated for {sys1_id} to {new_req_status}'})
    return jsonify({'status': 'error', 'message': f'SYS.1 Requirement with id {sys1_id} not found'}), 404

# Add endpoint to get traceability data (from session)
@app.route('/api/traceability_data')
def get_traceability_data():
    # Load both customer and SYS.1 requirements from session
    customer_reqs = session.get('customer_elicitation_requirements', [])
    sys1_reqs = session.get('sys1_elicitation_requirements', [])

    # For the traceability dashboard, we need a combined view
    # The simplest is to send both lists and let the frontend build the view
    return jsonify({
        'customer_requirements': customer_reqs,
        'sys1_requirements': sys1_reqs
    })

# New endpoint to update requirement details
@app.route('/api/agent1/update_requirement', methods=['POST'])
def update_requirement_details():
    try:
        data = request.get_json()
        if not data:
            return jsonify({'status': 'error', 'message': 'No data provided'}), 400

        sys1_id_to_update = data.get('sys1_id')
        updated_customer_req = data.get('customer_requirement')
        updated_sys1_req = data.get('sys1_requirement')
        # Add other fields here if they become editable in the modal

        if not sys1_id_to_update:
             return jsonify({'status': 'error', 'message': 'SYS.1 ID is required for update'}), 400

        sys1_requirements = session.get('sys1_elicitation_requirements', [])
        updated = False
        for req in sys1_requirements:
            if req.get('sys1_id') == sys1_id_to_update:
                # Update fields if provided
                if updated_customer_req is not None:
                    req['customer_requirement'] = updated_customer_req
                if updated_sys1_req is not None:
                    req['sys1_requirement'] = updated_sys1_req
                # Update other fields here

                session['sys1_elicitation_requirements'] = sys1_requirements # Save updated list back to session
                updated = True
                break

        if updated:
            # Also find the corresponding customer requirement and update it if needed
            # This assumes a 1:1 relationship or that the customer requirement text needs syncing
            # If 1:Many, this logic might need adjustment based on how customer_requirements are managed
            customer_requirements = session.get('customer_elicitation_requirements', [])
            for cust_req in customer_requirements:
                # Find the corresponding customer requirement by checking if its ID is in the sys1_req's customer_trace_ids
                # This handles both 1:1 and 1:Many traceability from the customer side to the SYS.1 side.
                # Note: This assumes sys1_req has a customer_trace_ids list
                if cust_req.get('customer_id') in req.get('customer_trace_ids', []):
                    if updated_customer_req is not None:
                        cust_req['customer_requirement'] = updated_customer_req
                    session['customer_elicitation_requirements'] = customer_requirements # Save updated list
                    break # Assuming one customer req per sys1 trace for simplicity here

            return jsonify({'status': 'success', 'message': f'Requirement {sys1_id_to_update} updated.'})
        else:
            return jsonify({'status': 'error', 'message': f'Requirement with SYS.1 ID {sys1_id_to_update} not found.'}), 404

    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

# New endpoint to handle feedback submission
@app.route('/api/submit_feedback', methods=['POST'])
def submit_feedback():
    data = request.json
    feedback = data.get('feedback')
    page = data.get('page', 'Unknown Page')
    print(f"Feedback from {page}: {feedback}")
    return jsonify({'status': 'success', 'message': 'Feedback received, thank you!'})

# Agent 2 Routes
@app.route('/api/agent2/process', methods=['POST'])
def process_sys1_for_agent2():
    print("[DEBUG] /api/agent2/process endpoint hit.") # Added log
    try:
        # Check if a file was uploaded
        if 'file' in request.files:
            print("[DEBUG] File upload detected.") # Added log
            file = request.files['file']
            if file.filename == '':
                print("[DEBUG] No selected file for upload.") # Added log
                return jsonify({'error': 'No selected file for upload'}), 400
            # Save the uploaded file temporarily
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filepath)
            input_data = {'file_path': filepath, 'source': 'upload'}
            print(f"[DEBUG] Processed file upload: {input_data}") # Added log

        # Check if processing the automatic file is requested
        elif request.form.get('source') == 'automatic_file':
            print("[DEBUG] Automatic file source requested.") # Added log
            # Define the automatic file path
            root_path = os.path.dirname(os.path.abspath(__file__))
            automatic_file_path = os.path.join(root_path, 'Inputs', 'sys1_requirements.xlsx')
            print(f"[DEBUG] Expected automatic file path: {automatic_file_path}") # Added log
            if not os.path.exists(automatic_file_path):
                 print("[DEBUG] Automatic file not found.") # Added log
                 return jsonify({'status': 'error', 'message': f'Automatic file not found at {automatic_file_path}'}), 404
            input_data = {'file_path': automatic_file_path, 'source': 'automatic'}
            print(f"[DEBUG] Processed automatic file source: {input_data}") # Added log

        # Check for raw text input
        elif request.form.get('raw_content'):
            print("[DEBUG] Raw text input detected.") # Added log
            raw_content = request.form.get('raw_content')
            input_data = {'raw_content': raw_content, 'source': 'raw_text'}
            print(f"[DEBUG] Processed raw text input: {input_data}") # Added log

        # Check for session source (used by manual input after upload)
        elif request.form.get('source') == 'session':
             print("[DEBUG] Session source requested.") # Added log
             # In this case, Sys2Agent should read from the session
             input_data = {'source': 'session'}
             print(f"[DEBUG] Processed session source: {input_data}") # Added log

        else:
            print("[DEBUG] No valid input source provided.") # Added log
            return jsonify({'status': 'error', 'message': 'No input provided (file, automatic_file request, raw_content, or session)'}), 400

        # Initialize and run Agent 2
        print("[DEBUG] Initializing Sys2Agent.") # Added log
        process_result = sys2_agent.process_sys1_input(input_data)

        # Assuming process_result contains sys2_requirements, dependencies, etc.
        print(f"[DEBUG] Sys2Agent process_sys1_input returned: {process_result.get('status')}") # Added log
        if process_result and process_result.get('status') == 'success':
             # Store results in session
             session['sys2_requirements'] = process_result.get('sys2_requirements', [])
             session['dependencies'] = process_result.get('dependencies', [])
             # Store other results as needed (e.g., classification, verification mapping)
             session['sys2_classification'] = process_result.get('classification', {})
             session['sys2_verification_mapping'] = process_result.get('verification_mapping', {})

             print("[DEBUG] SYS.2 data stored in session.") # Added log
             return jsonify({
                 'status': 'success',
                 'message': 'SYS.2 requirements and data generated.',
                 'sys2_requirements': session['sys2_requirements'],
                 'dependencies': session['dependencies'],
                 'classification': session['sys2_classification'],
                 'verification_mapping': session['sys2_verification_mapping']
             })
        else:
             print("[DEBUG] Agent 2 processing failed.") # Added log
             return jsonify({'status': 'error', 'message': process_result.get('message', 'Agent 2 processing failed.')}), 500

    except Exception as e:
        print(f"[DEBUG] Error in /api/agent2/process endpoint: {e}") # Added log
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/agent2/update_requirement', methods=['POST'])
def update_sys2_requirement():
    try:
        data = request.get_json()
        if not data or 'id' not in data or 'updates' not in data:
            return jsonify({'status': 'error', 'message': 'Missing required fields: id or updates'}), 400

        sys2_id_to_update = data.get('id')
        updates = data.get('updates')

        # Load SYS.2 requirements from session
        sys2_requirements = session.get('sys2_requirements', [])
        updated = False

        # Find and update the requirement
        for req in sys2_requirements:
            if req.get('sys2_id') == sys2_id_to_update:
                for field, value in updates.items():
                    req[field] = value
                updated = True
                break

        if updated:
            # Save the updated list back to session
            session['sys2_requirements'] = sys2_requirements
            session.modified = True # Mark session as modified
            return jsonify({'status': 'success', 'message': f'Requirement {sys2_id_to_update} updated successfully.'})
        else:
            return jsonify({'status': 'error', 'message': f'Requirement with SYS.2 ID {sys2_id_to_update} not found.'}), 404

    except Exception as e:
        print(f"Error updating SYS.2 requirement: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/agent2/export/<format>', methods=['GET'])
def export_sys2_requirements(format):
    try:
        sys2_requirements = session.get('sys2_requirements', [])
        if not sys2_requirements:
            return jsonify({'status': 'error', 'message': 'No SYS.2 requirements to export.'}), 400

        # Only export SYS.2 Req. ID and SYS.2 System Requirement for XLSX
        if format == 'xlsx':
            export_data_io = sys2_agent.export_requirements(sys2_requirements, 'xlsx', export_fields_list=['sys2_id', 'sys2_requirement'])
            if export_data_io is None:
                return jsonify({'status': 'error', 'message': 'Failed to generate export data.'}), 500
            return send_file(export_data_io, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', as_attachment=True, download_name='sys2_requirements.xlsx')
        else:
            # For other formats, use the default export fields
            export_data = sys2_agent.export_requirements(sys2_requirements, format)
            if export_data is None:
                return jsonify({'status': 'error', 'message': 'Failed to generate export data.'}), 500
            if format == 'csv':
                return send_file(io.BytesIO(export_data.encode('utf-8')), mimetype='text/csv', as_attachment=True, download_name='sys2_requirements.csv')
            elif format == 'txt':
                return send_file(io.BytesIO(export_data.encode('utf-8')), mimetype='text/plain', as_attachment=True, download_name='sys2_requirements.txt')
            elif format == 'docx':
                return send_file(export_data, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document', as_attachment=True, download_name='sys2_requirements.docx')
            elif format == 'pdf':
                return send_file(export_data, mimetype='application/pdf', as_attachment=True, download_name='sys2_requirements.pdf')
            else:
                return jsonify({'status': 'error', 'message': f'Unsupported export format: {format}'}), 400
    except Exception as e:
        print(f"[ERROR] Error exporting SYS.2 requirements: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': f'Error exporting SYS.2 requirements: {str(e)}'}), 500

@app.route('/api/agent2/combined_dashboard_summary', methods=['GET'])
def get_agent2_combined_dashboard_summary():
    """
    Endpoint to provide combined summary statistics for SYS.1 and SYS.2 requirements for the Agent 2 dashboard.
    """
    try:
        # Load SYS.1 and SYS.2 requirements from session
        sys1_requirements = session.get('sys1_elicitation_requirements', [])
        sys2_requirements = session.get('sys2_requirements', [])

        # Add print statement to check if requirements are loaded
        print(f"[DEBUG] Combined dashboard summary request. SYS.1 requirements in session: {len(sys1_requirements) if sys1_requirements else 0}")
        print(f"[DEBUG] Combined dashboard summary request. SYS.2 requirements in session: {len(sys2_requirements) if sys2_requirements else 0}")

        # Calculate SYS.1 summary (logic adapted from Agent 1 export endpoint)
        customer_requirements = session.get('customer_elicitation_requirements', [])
        customer_map = {req.get('customer_id'): req for req in customer_requirements if req.get('customer_id')}

        total_customer_reqs = len(customer_requirements)
        traced_customer_ids = set()
        for sys1_req in sys1_requirements:
            if sys1_req.get('customer_trace_ids'):
                for cust_id in sys1_req['customer_trace_ids']:
                    if cust_id in customer_map:
                        traced_customer_ids.add(cust_id)
        traced_customer_count = len(traced_customer_ids)
        untraced_customer_count = total_customer_reqs - traced_customer_count

        total_sys1_reqs = len(sys1_requirements)
        approved_sys1_count = sum(1 for req in sys1_requirements if req.get('req_status') == 'Approved')
        rejected_sys1_count = sum(1 for req in sys1_requirements if req.get('req_status') == 'Rejected')
        draft_sys1_count = total_sys1_reqs - approved_sys1_count - rejected_sys1_count

        sys1_summary = {
            'total_sys1_reqs': total_sys1_reqs,
            'traced_to_customer': traced_customer_count,
            'not_traced_to_customer': untraced_customer_count,
            'status_breakdown': {
                'Draft': draft_sys1_count,
                'Reviewed': 0, # SYS.1 doesn't have Reviewed status in current logic
                'Approved': approved_sys1_count,
                'Rejected': rejected_sys1_count
            }
        }

        # Calculate SYS.2 summary using Sys2Agent method
        sys2_summary = sys2_agent.get_dashboard_summary(sys2_requirements)

        # Include the full requirements lists in the summary response for the traceability table
        return jsonify({
            'status': 'success',
            'summary': {
                'sys1': sys1_summary,
                'sys2': sys2_summary,
                'sys1_requirements': sys1_requirements,
                'sys2_requirements': sys2_requirements
            }
        })

    except Exception as e:
        print(f"Error fetching combined dashboard summary: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

# Agent 3 Routes (Placeholder)
@app.route('/api/agent3/review', methods=['POST'])
def review_sys2_requirements():
    # In a real application, you might get selected SYS.2 IDs from the request body
    # For now, let's process all SYS.2 requirements currently in the session
    sys2_requirements = session.get('sys2_requirements', [])

    if not sys2_requirements:
        return jsonify({'status': 'error', 'message': 'No SYS.2 requirements found in session.'}), 400

    agent3 = ReviewAgent()
    review_results, suggestions, test_proposals = agent3.review_requirements(sys2_requirements)

    # You might want to store the review results, suggestions, and test proposals in the session
    session['review_results'] = review_results
    session['suggestions'] = suggestions
    session['test_proposals'] = test_proposals

    return jsonify({
        'status': 'success',
        'message': 'SYS.2 requirements reviewed and suggestions/test proposals generated.',
        'review_results': review_results,
        'suggestions': suggestions,
        'test_proposals': test_proposals
    })

@app.route('/api/agent3/process_sys2', methods=['POST'])
def process_sys2():
    print("[DEBUG] /api/agent3/process_sys2 endpoint hit.")
    try:
        # Check if a file was uploaded
        if 'file' in request.files:
            print("[DEBUG] File upload detected for Agent 3.")
            file = request.files['file']
            if file.filename == '':
                print("[DEBUG] No selected file for upload for Agent 3.")
                return jsonify({'error': 'No selected file for upload'}), 400
            # Save the uploaded file temporarily
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filepath)
            input_data = {'file_path': filepath, 'source': 'upload'}
            print(f"[DEBUG] Processed file upload for Agent 3: {input_data}")

        # Check if processing the automatic file is requested
        elif request.form.get('source') == 'automatic_file':
            print("[DEBUG] Automatic file source requested for Agent 3.")
            # Define the automatic file path using the user's specified path
            automatic_file_path = r'D:\AgentX\AutoTestGen_MAPS_Agents123\AutoTestGen_MAPS\Inputs\sys2_requirements.xlsx'
            print(f"[DEBUG] Expected automatic file path for Agent 3: {automatic_file_path}")
            if not os.path.exists(automatic_file_path):
                 print("[DEBUG] Automatic file not found for Agent 3.")
                 return jsonify({'status': 'error', 'message': f'Automatic file not found at {automatic_file_path}'}), 404
            input_data = {'file_path': automatic_file_path, 'source': 'automatic'}
            print(f"[DEBUG] Processed automatic file source for Agent 3: {input_data}")

        # Add other potential input sources if needed (e.g., raw text, session)
        # elif request.form.get('raw_content'):
        #     print("[DEBUG] Raw text input detected for Agent 3.")
        #     raw_content = request.form.get('raw_content')
        #     input_data = {'raw_content': raw_content, 'source': 'raw_text'}
        #     print(f"[DEBUG] Processed raw text input for Agent 3: {input_data}")
        # elif request.form.get('source') == 'session':
        #      print("[DEBUG] Session source requested for Agent 3.")
        #      input_data = {'source': 'session'}
        #      print(f"[DEBUG] Processed session source for Agent 3: {input_data}")

        else:
            print("[DEBUG] No valid input source provided for Agent 3.")
            return jsonify({'status': 'error', 'message': 'No input provided (file or automatic_file request for Agent 3)'}), 400

        # Initialize and run Agent 3 (ReviewAgent)
        print("[DEBUG] Initializing ReviewAgent for Agent 3.")
        agent3 = ReviewAgent()

        requirements_to_process = []
        processing_source = input_data.get('source', 'unknown')

        try:
            if processing_source == 'upload' and input_data.get('file_path'):
                # Read requirements directly from the uploaded file within the endpoint
                file_path = input_data.get('file_path')
                print(f"[DEBUG] Reading requirements from uploaded file: {file_path}")
                # Pass the file_path to the agent's method to handle reading
                requirements_to_process = agent3._read_requirements_from_excel(file_path)
                print(f"[DEBUG] Read {len(requirements_to_process)} requirements from uploaded file.")

            elif processing_source == 'automatic_file' and input_data.get('file_path'):
                 # Read requirements from the automatic file within the endpoint
                 file_path = input_data.get('file_path')
                 print(f"[DEBUG] Attempting to read automatic file: {file_path}")

                 # Explicitly check if the file exists and is readable
                 if not os.path.exists(file_path):
                     print(f"[ERROR] Automatic file not found at {file_path}")
                     return jsonify({'status': 'error', 'message': f'Automatic file not found at {file_path}'}), 404
                 if not os.access(file_path, os.R_OK):
                      print(f"[ERROR] Automatic file not readable at {file_path}")
                      return jsonify({'status': 'error', 'message': f'Automatic file not readable at {file_path}. Check permissions.'}), 500 # Use 500 for permission errors

                 print(f"[DEBUG] Automatic file found and is readable: {file_path}")
                 # Pass the file_path to the agent's method to handle reading
                 requirements_to_process = agent3._read_requirements_from_excel(file_path)
                 print(f"[DEBUG] Read {len(requirements_to_process)} requirements from automatic file.")

            # Add other sources if needed (e.g., session data)

            if not requirements_to_process:
                # If no requirements were read, return an error before processing
                print("[DEBUG] No requirements read from input source.")
                return jsonify({'status': 'error', 'message': 'Could not read requirements from the provided input.'}), 400

            # Prepare input data for the ReviewAgent.process method
            agent_process_input = {
                'requirements': requirements_to_process,
                'source': processing_source # Pass the original source
                # Add other necessary input fields for the agent if required by agent3.process
            }

            print("[DEBUG] Calling ReviewAgent process with", len(requirements_to_process), "requirements.")
            # Call the actual ReviewAgent process method with the prepared input
            process_result = agent3.process(agent_process_input)

            if process_result and process_result.get('status') == 'success':
                 # Store results in session (using different keys for Agent 3)
                 # Ensure these keys match what the frontend expects
                 session['agent3_sys2_requirements'] = process_result.get('sys2_requirements_for_review', [])
                 session['agent3_compliance_results'] = process_result.get('compliance_results', [])
                 session['agent3_suggestions'] = process_result.get('suggestions', [])

                 print("[DEBUG] Agent 3 data stored in session.")
                 return jsonify({
                     'status': 'success',
                     'message': process_result.get('message', 'Agent 3 processing successful!'), # Use message from agent process if available
                     'sys2_requirements_for_review': session['agent3_sys2_requirements'],
                     'compliance_results': session['agent3_compliance_results'],
                     'suggestions': session['agent3_suggestions']
                 })
            else:
                 print("[DEBUG] Agent 3 processing failed in agent.process method.")
                 # Return the specific error message from the agent's result if available
                 error_message = process_result.get('message', 'Agent 3 processing failed.')
                 return jsonify({'status': 'error', 'message': error_message}), 500

        except FileNotFoundError as fnf_error:
            print(f"[ERROR] File not found during Agent 3 processing: {fnf_error}")
            traceback.print_exc()
            return jsonify({'status': 'error', 'message': f'Error: Input file not found. {str(fnf_error)}'}), 404
        except Exception as e:
            print(f"[ERROR] An unexpected error occurred during Agent 3 processing: {e}")
            import traceback
            traceback.print_exc()
            return jsonify({'status': 'error', 'message': f'An error occurred during processing: {str(e)}'}), 500

    except Exception as e:
        print(f"[ERROR] An error occurred in /api/agent3/process_sys2 endpoint: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': f'An error occurred in the endpoint: {str(e)}'}), 500

@app.route('/api/agent3/review/<sys2_id>', methods=['GET'])
def get_agent3_feedback(sys2_id):
    try:
        # Get SYS.2 requirements from session
        sys2_requirements = session.get('sys2_requirements', [])
        
        # Find the specific requirement
        requirement = next((req for req in sys2_requirements if req.get('sys2_id') == sys2_id), None)
        
        if not requirement:
            return jsonify({'status': 'error', 'message': f'Requirement {sys2_id} not found'}), 404

        # Initialize review agent
        review_agent = ReviewAgent()
        
        # Get review feedback for this requirement
        feedback = review_agent.get_requirement_feedback(requirement)
        
        return jsonify({
            'status': 'success',
            'review_status': feedback.get('status', 'Pending'),
            'suggestions': feedback.get('suggestions', []),
            'compliance_results': feedback.get('compliance_results', []),
            'test_proposals': feedback.get('test_proposals', [])
        })
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/agent2/export/sys2_specific_xlsx', methods=['POST'])
def export_sys2_specific_xlsx():
    """
    Exports only SYS.2 Req. ID and SYS.2 System Requirement to a specific XLSX file path.
    """
    try:
        # Load SYS.2 requirements from session
        sys2_requirements = session.get('sys2_requirements', [])

        if not sys2_requirements:
            return jsonify({'status': 'error', 'message': 'No SYS.2 requirements found in session to export.'}), 400

        # Define the specific fields to export
        fields_to_export = ['sys2_id', 'sys2_requirement']

        # Call the export method with the specific fields and format
        export_data_io = sys2_agent.export_requirements(sys2_requirements, 'xlsx', export_fields_list=fields_to_export)

        if export_data_io is None:
             return jsonify({'status': 'error', 'message': 'Failed to generate export data.'}), 500

        # Define the absolute path to save the file (updated to D: drive as requested)
        root_path = os.path.dirname(os.path.abspath(__file__))
        output_path = os.path.join(root_path, 'Inputs', 'sys2_requirements.xlsx')

        # Ensure the directory exists
        output_dir = os.path.dirname(output_path)
        os.makedirs(output_dir, exist_ok=True)

        # Save the BytesIO content to the specified file path
        with open(output_path, 'wb') as f:
            f.write(export_data_io.getvalue())

        print(f"[INFO] Successfully exported specific SYS.2 requirements to {output_path}")

        return jsonify({
            'status': 'success',
            'message': f'Successfully exported SYS.2 requirements (ID, Requirement) to {output_path}',
            'file_path': output_path
        })

    except Exception as e:
        print(f"[ERROR] Error exporting specific SYS.2 requirements: {e}")
        return jsonify({'status': 'error', 'message': f'Error exporting SYS.2 requirements: {str(e)}'}), 500

@app.route('/api/agent2/save_sys2_xlsx', methods=['POST'])
def save_sys2_xlsx():
    """
    Endpoint to save SYS.2 requirements as an XLSX file to a specific local path.
    Triggered automatically after successful Agent 2 processing.
    """
    try:
        # Load SYS.2 requirements from session
        sys2_requirements = session.get('sys2_requirements', [])

        if not sys2_requirements:
            print("[DEBUG] save_sys2_xlsx: No SYS.2 requirements found in session to save.")
            return jsonify({'status': 'error', 'message': 'No SYS.2 requirements found in session to save.'}), 400

        # Define the absolute path to save the file
        # Use the workspace path from user_info if available, otherwise hardcode
        # Using the hardcoded path provided by the user for now.
        root_path = os.path.dirname(os.path.abspath(__file__))
        output_path = os.path.join(root_path, 'Inputs', 'sys2_requirements.xlsx')

        # Ensure the directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # Define the fields to export (using the default dashboard export fields)
        # Based on Sys2Agent.export_requirements default fields
        fields_to_export = [
            'sys1_id', 'sys1_requirement',
            'sys2_id', 'sys2_requirement',
            'classification',
            'verification_mapping',
            'verification_criteria',
            'domain', 'priority', 'rationale', 'req_status'
        ]

        # Call the export method to get the BytesIO object
        # Passing None for export_fields_list uses the default fields defined inside the method
        # Or we can explicitly pass fields_to_export defined above
        # export_data_io = sys2_agent.export_requirements(sys2_requirements, 'xlsx', export_fields_list=None) # Use default fields
        # Pass the specific fields for the automatic export
        export_data_io = sys2_agent.export_requirements(sys2_requirements, 'xlsx', export_fields_list=['sys2_id', 'sys2_requirement'])

        if export_data_io is None:
            print("[DEBUG] save_sys2_xlsx: Failed to generate export data (export_data_io is None).")
            return jsonify({'status': 'error', 'message': 'Failed to generate export data.'}), 500

        # Save the BytesIO content to the specified file path
        with open(output_path, 'wb') as f:
            f.write(export_data_io.getvalue())

        print(f"[INFO] Successfully saved SYS.2 requirements to {output_path}")

        return jsonify({
            'status': 'success',
            'message': f'SYS.2 requirements saved to {output_path}',
            'file_path': output_path # Optionally return the file path
        })

    except Exception as e:
        print(f"[ERROR] Error saving SYS.2 requirements to XLSX: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': f'Error saving SYS.2 requirements to XLSX: {str(e)}'}), 500

@app.route('/api/agent3/export_accepted', methods=['POST'])
def export_accepted_requirements():
    try:
        data = request.get_json()
        requirements = data.get('requirements', [])
        
        if not requirements:
            return jsonify({
                'status': 'error',
                'message': 'No requirements provided for export'
            }), 400

        # Export the requirements
        review_agent._export_accepted_requirements(requirements)
        
        return jsonify({
            'status': 'success',
            'export_message': f'Successfully exported {len(requirements)} accepted requirements to sys2_requirements_reviewed.xlsx'
        })
        
    except Exception as e:
        return jsonify({
            'status': 'error',
            'message': str(e)
        }), 500

@app.route('/agent4/dashboard_data', methods=['GET'])
def get_dashboard_data():
    """Provide data for Agent 4 dashboard charts and tables"""
    try:
        # Retrieve data from session
        test_cases = session.get('agent4_test_cases', [])
        traceability_matrix = session.get('agent4_traceability_matrix', {})
        coverage_analysis = session.get('agent4_coverage_analysis', {})
        maturity_status = session.get('agent4_maturity_status', {})
        requirements = session.get('agent4_requirements', [])
        
        if not test_cases and not requirements:
            return jsonify({'status': 'info', 'message': 'No data available. Generate test cases first.'}), 200
        
        return jsonify({
            'status': 'success',
            'test_cases': test_cases,
            'traceability_matrix': traceability_matrix,
            'coverage_analysis': coverage_analysis,
            'maturity_status': maturity_status,
            'requirements': requirements
        })
    except Exception as e:
        return jsonify({
            'status': 'error',
            'message': str(e)
        }), 500

@app.route('/agent4/update_test_case', methods=['POST'])
def update_test_case():
    """Handle updates to individual test cases"""
    try:
        data = request.get_json();
        test_case_id_to_update = data.get('test_case_id')
        updates = data.get('updates')

        if not test_case_id_to_update or not updates:
            return jsonify({'status': 'error', 'message': 'Missing test_case_id or updates'}), 400

        # Load test cases from session
        test_cases = session.get('agent4_test_cases', [])
        updated_test_case = None
        
        # Find and update the test case
        for tc in test_cases:
            if tc.get('test_id') == test_case_id_to_update:
                for field, value in updates.items():
                    # Handle specific fields that might need list conversion (Steps, Expected Results)
                    if field in ['steps', 'expected_results'] and isinstance(value, str):
                         # Convert newline-separated string back to list
                        tc[field] = [item.strip() for item in value.split('\n') if item.strip()]
                    else:
                        tc[field] = value
                updated_test_case = tc
                break

        if updated_test_case:
            # Save the updated list back to session
            session['agent4_test_cases'] = test_cases
            session.modified = True # Mark session as modified
            
            return jsonify({
                'status': 'success',
                'message': f'Test case {test_case_id_to_update} updated successfully.',
                'updated_test_case': updated_test_case # Return the updated test case
            })
        else:
            return jsonify({'status': 'error', 'message': f'Test case with ID {test_case_id_to_update} not found.'}), 404

    except Exception as e:
        print(f"Error updating test case: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/agent4/submit_correction', methods=['POST'])
def submit_correction():
    """Receive and store user corrections for test cases"""
    try:
        data = request.get_json()
        test_case_id = data.get('test_case_id')
        corrected_data = data.get('corrected_data')

        if not test_case_id or not corrected_data:
            return jsonify({'status': 'error', 'message': 'Missing test_case_id or corrected_data'}), 400

        # Load existing corrections from session, or initialize if none exist
        corrections = session.get('agent4_corrections', {})

        # Store the corrected data for the specific test case ID
        corrections[test_case_id] = corrected_data

        # Save the updated corrections back to session
        session['agent4_corrections'] = corrections
        session.modified = True # Mark session as modified

        print(f"[INFO] Received and stored correction for test case: {test_case_id}")
        # print(f"[DEBUG] Stored corrections: {corrections}") # Optional: log full corrections

        return jsonify({
            'status': 'success',
            'message': f'Correction for test case {test_case_id} submitted successfully.'
        })

    except Exception as e:
        print(f"[ERROR] Error submitting correction: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/agent4/load_sys2', methods=['GET'])
def agent4_load_sys2():
    """Auto-load sys2_requirements_reviewed.xlsx and generate test cases."""
    try:
        file_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Inputs', 'sys2_requirements_reviewed.xlsx')
        requirements = testgen_agent.load_requirements(file_path)
        test_cases = testgen_agent.generate_test_cases(requirements)
        session['agent4_requirements'] = requirements
        session['agent4_test_cases'] = test_cases
        session.modified = True
        return jsonify({'status': 'success', 'test_cases': test_cases})
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)})

@app.route('/api/agent4/upload_sys2', methods=['POST'])
def agent4_upload_sys2():
    """Handle manual upload of SYS.2 file and generate test cases."""
    try:
        file = request.files.get('file')
        if not file or file.filename == '':
            return jsonify({'status': 'error', 'message': 'No file uploaded.'}), 400
        # Save to a temp location
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(temp_path)
        requirements = testgen_agent.load_requirements(temp_path)
        test_cases = testgen_agent.generate_test_cases(requirements)
        session['agent4_requirements'] = requirements
        session['agent4_test_cases'] = test_cases
        session.modified = True
        os.remove(temp_path)
        return jsonify({'status': 'success', 'test_cases': test_cases})
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)})

@app.route('/api/agent4/test_cases', methods=['GET'])
def agent4_get_test_cases():
    """Return generated test cases as JSON."""
    test_cases = session.get('agent4_test_cases', [])
    return jsonify({'status': 'success', 'test_cases': test_cases})

def open_browser():
    time.sleep(1) # Give the server a moment to start
    webbrowser.open('http://localhost:5000/')

if __name__ == '__main__':
    # Open browser in a separate thread after a short delay
    threading.Thread(target=open_browser).start()
    app.run(debug=True) 